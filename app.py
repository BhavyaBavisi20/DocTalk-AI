import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from io import BytesIO
from langchain_community.vectorstores import FAISS
from langchain.text_splitter.recursive import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LCDocument
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from gtts import gTTS
import base64
import asyncio

# --- Page Configuration and CSS ---
st.set_page_config(
    page_title="DocTalk AI",
    page_icon="✨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for a cleaner, Gemini-like UI
st.markdown("""
    <style>
:root {
    --chat-text-color: #333;
    --chat-bg-color: #ffffff;
    --user-bg-color: #e6e6fa;
    --assistant-bg-color: #f0f8ff;
}

@media (prefers-color-scheme: dark) {
    :root {
        --chat-text-color: #f0f0f0;
        --chat-bg-color: #1e1e1e;
        --user-bg-color: #4b3f72;
        --assistant-bg-color: #2c3e50;
    }
}

.stChatMessage {
    background-color: var(--chat-bg-color) !important;
    border-radius: 15px;
    padding: 15px 20px;
    margin-bottom: 10px;
    box-shadow: 0 2px 5px rgba(0,0,0,0.05);
    color: var(--chat-text-color) !important;
}

.stChatMessage[data-testid="stChatMessage-user"] {
    background-color: var(--user-bg-color) !important;
    align-self: flex-end;
    text-align: right;
    border-bottom-right-radius: 5px;
    color: var(--chat-text-color) !important;
}

.stChatMessage[data-testid="stChatMessage-assistant"] {
    background-color: var(--assistant-bg-color) !important;
    align-self: flex-start;
    text-align: left;
    border-bottom-left-radius: 5px;
    color: var(--chat-text-color) !important;
}

.stChatMessage p,
.stChatMessage span,
.stChatMessage div {
    color: var(--chat-text-color) !important;
}

/* Optional: style chat input too */
.stTextInput > div > div > input {
    border-radius: 25px;
    padding: 10px 15px;
    border: 1px solid #ddd;
    color: var(--chat-text-color) !important;
    background-color: var(--chat-bg-color) !important;
}
</style>

""", unsafe_allow_html=True)

# --- Configure Gemini ---
# Ensure GEMINI_API_KEY is set in Streamlit secrets
try:
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    model = genai.GenerativeModel("gemini-2.5-flash")
except Exception as e:
    st.error(f"Configuration Error: Please ensure 'GEMINI_API_KEY' is set in your Streamlit secrets.toml file. Details: {e}")
    st.stop()
# models = genai.list_models()
# for m in models:
#     st.write(m.name)


# --- Session State Setup ---
if "chat_history" not in st.session_state:
    st.session_state.chat_history = [] # Stores (user_message, assistant_response, audio_base64)

if "documents" not in st.session_state:
    st.session_state.documents = [] # Stores extracted text from documents

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None # FAISS vector store for document embeddings

if "show_file_uploader" not in st.session_state:
    st.session_state.show_file_uploader = False # Controls visibility of in-chat file uploader

if "voice_output_enabled" not in st.session_state:
    st.session_state.voice_output_enabled = False # Controls voice output

# --- Helper Functions ---

# Function to extract text from various document types
def extract_text(file):
    """Extracts text content from uploaded files based on their type."""
    name = file.name.lower()
    try:
        if name.endswith(".pdf"):
            reader = PdfReader(file)
            return "\n".join(p.extract_text() for p in reader.pages if p.extract_text())
        elif name.endswith(".docx"):
            doc = Document(file)
            return "\n".join(p.text for p in doc.paragraphs)
        elif name.endswith(".pptx"):
            ppt = Presentation(file)
            # Iterate through slides and shapes to extract text
            return "\n".join(shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif name.endswith(".xlsx"):
            wb = openpyxl.load_workbook(file, data_only=True)
            out = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    # Join cells in a row with tab, handle None values
                    out.append("\t".join(str(cell) if cell is not None else "" for cell in row))
            return "\n".join(out)
        elif name.endswith((".txt", ".md")):
            return file.read().decode("utf-8")
        else:
            return "" # Return empty string for unsupported types
    except Exception as e:
        # Log and return an error message if extraction fails
        st.error(f"❌ Error reading {file.name}: {e}")
        return ""

# Function to get the embedding model
@st.cache_resource
def get_embedding_model():
    """Initializes and caches the Google Generative AI Embeddings model."""
    # Ensure an event loop is available for async operations if running in a thread
    try:
        asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    return GoogleGenerativeAIEmbeddings(
        model="models/embedding-001",
        google_api_key=st.secrets["GEMINI_API_KEY"]
    )

# Function to update the FAISS vector store with new documents
from langchain.text_splitter import RecursiveCharacterTextSplitter

def update_vectorstore():
    """
    Updates the FAISS vector store with text from st.session_state.documents.
    Uses safe chunking to avoid GoogleGenerativeAI embedding failures.
    """

    # Safe, Gemini-compatible splitter
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1500,        # Under embedding-001 token limit
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", " ", ""]
    )

    docs = []

    for doc_text in st.session_state.documents:

        # Clean malformed Unicode (common in PDFs)
        cleaned = doc_text.encode("utf-8", "ignore").decode()

        # Split into safe-size chunks
        pieces = text_splitter.split_text(cleaned)

        # Build LangChain Documents
        for p in pieces:
            if p.strip():
                docs.append(LCDocument(page_content=p))

    # No documents after cleaning → no vector store
    if not docs:
        st.session_state.vector_store = None
        return

    embedding_model = get_embedding_model()

    # Build or update FAISS store
    try:
        if st.session_state.vector_store:
            st.session_state.vector_store.add_documents(docs)
        else:
            st.session_state.vector_store = FAISS.from_documents(docs, embedding_model)

        st.success("✅ Documents processed and ready for chat!")

    except Exception as e:
        st.error(f"❌ Vector store update failed: {e}")

# Function to ask a question using RAG (Retrieval Augmented Generation)
def ask_question_vector(query):
    """
    Retrieves relevant document snippets from the vector store and
    uses them as context for the Gemini model to answer the query.
    """
    if not st.session_state.vector_store:
        return "⚠️ Please upload a document first to enable document-based questioning."

    try:
        # Perform similarity search to get top relevant document chunks
        results = st.session_state.vector_store.similarity_search(query, k=5)
        context = "\n---\n".join([doc.page_content for doc in results])

        # Construct the prompt for the Gemini model
        prompt = f"""
        You are a helpful AI assistant.
        Use the following document snippets to answer the user's question.
        If the question cannot be answered from the provided document snippets,
        state that you don't have enough information in the documents.

        Document Snippets:
        {context}

        Question: {query}
        """
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        st.error(f"An error occurred during response generation: {e}")
        return "An error occurred while trying to answer your question. Please try again."

# Function to convert text to speech and return as base64 audio
def text_to_audio_base64(text):
    """Converts text to speech using gTTS and returns base64 encoded audio."""
    try:
        tts = gTTS(text=text, lang='en', slow=False)
        audio_bytes_io = BytesIO()
        tts.write_to_fp(audio_bytes_io)
        audio_bytes_io.seek(0) # Rewind to the beginning of the BytesIO object
        # Encode audio bytes to base64 for embedding in HTML audio tag
        return base64.b64encode(audio_bytes_io.read()).decode('utf-8')
    except Exception as e:
        st.error(f"Error generating audio: {e}")
        return None

# Function to format chat history for download
def get_chat_history_for_download():
    """Formats the entire chat history into a readable string for download."""
    history_text = ""
    for q, a, _ in st.session_state.chat_history: # _ is for audio_base64
        history_text += f"User: {q}\n"
        history_text += f"Assistant: {a}\n\n"
    return history_text

# --- UI Layout ---

st.markdown("<h2 style='color:#a678f2; text-align: center;'>DocTalk AI ✨</h2>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center;'>Chat with your documents, powered by Gemini!</p>", unsafe_allow_html=True)

# Sidebar for document management and conversation download
with st.sidebar:
    st.markdown("### 📎 Document Management")
    st.markdown("Upload documents here to make them available for chat.")
    sidebar_uploaded_files = st.file_uploader(
        "Supported: PDF, DOCX, PPTX, XLSX, TXT, MD",
        type=["pdf", "docx", "pptx", "xlsx", "txt", "md"],
        accept_multiple_files=True,
        key="sidebar_uploader" # Unique key for this uploader
    )
    if sidebar_uploaded_files:
        for file in sidebar_uploaded_files:
            text = extract_text(file)
            if text:
                st.session_state.documents.append(text)
        update_vectorstore()

    st.markdown("---")
    st.markdown("### ⚙️ Settings")
    # Toggle for voice output
    st.session_state.voice_output_enabled = st.checkbox(
        "Enable Voice Output",
        value=st.session_state.voice_output_enabled,
        help="Toggle to enable or disable the assistant's voice responses."
    )

    st.markdown("---")
    st.markdown("### ⬇️ Conversation Options")
    # Download conversation button
    if st.session_state.chat_history:
        st.download_button(
            label="Download Chat History",
            data=get_chat_history_for_download(),
            file_name="doctalk_ai_chat_history.txt",
            mime="text/plain"
        )
    else:
        st.info("Start a conversation to enable download.")

# Main chat display area
chat_placeholder = st.container()

with chat_placeholder:
    for i, (q, a, audio_b64) in enumerate(st.session_state.chat_history):
        # User message
        with st.chat_message("user", avatar="👤"):
            st.markdown(q)
        # Assistant message
        with st.chat_message("assistant", avatar="🤖"):
            st.markdown(a)
            # Only display audio if voice output is enabled and audio exists
            if st.session_state.voice_output_enabled and audio_b64:
                # Embed audio player
                st.markdown(f'<audio controls autoplay style="width: 100%;"><source src="data:audio/mpeg;base64,{audio_b64}" type="audio/mpeg"></audio>', unsafe_allow_html=True)


# --- Chat Input and Action Buttons ---
# Persistent input area at the bottom
st.markdown("---") # Separator for input area
input_col, upload_col, voice_col = st.columns([8, 1, 1])

with input_col:
    prompt = st.chat_input("Ask anything about your document...", key="chat_input_main")

with upload_col:
    # Button to toggle file uploader visibility in chat
    if st.button("📎 Attach", key="attach_button"):
        st.session_state.show_file_uploader = not st.session_state.show_file_uploader
        # Rerun to show/hide uploader immediately
        st.rerun()

with voice_col:
    # Placeholder for voice input (explains limitations)
    if st.button("🎤 Voice", key="voice_button"):
        st.warning("Voice input (Speech-to-Text) is not directly supported in Streamlit's frontend without custom components or browser-side APIs. Please type your message.")


# Conditional file uploader in chat area
if st.session_state.show_file_uploader:
    with st.container():
        st.markdown("---")
        st.markdown("### Upload documents directly here:")
        in_chat_uploaded_files = st.file_uploader(
            "Select files",
            type=["pdf", "docx", "pptx", "xlsx", "txt", "md"],
            accept_multiple_files=True,
            key="in_chat_uploader" # Unique key for this uploader
        )
        if in_chat_uploaded_files:
            for file in in_chat_uploaded_files:
                text = extract_text(file)
                if text:
                    st.session_state.documents.append(text)
            update_vectorstore()
            st.session_state.show_file_uploader = False # Hide after upload
            st.rerun() # Rerun to update chat and hide uploader


# Process user query
if prompt:
    # Add user message to chat history
    st.session_state.chat_history.append((prompt, "", None)) # User message, empty assistant response, no audio yet

    with chat_placeholder:
        with st.chat_message("user", avatar="👤"):
            st.markdown(prompt)

        with st.chat_message("assistant", avatar="🤖"):
            with st.spinner("🤖 Thinking..."):
                answer = ask_question_vector(prompt)
                st.markdown(answer)
                audio_b64 = None
                # Generate audio only if voice output is enabled
                if st.session_state.voice_output_enabled:
                    audio_b64 = text_to_audio_base64(answer)
                    if audio_b64:
                        st.markdown(f'<audio controls autoplay style="width: 100%;"><source src="data:audio/mpeg;base64,{audio_b64}" type="audio/mpeg"></audio>', unsafe_allow_html=True)

            # Update the last assistant response and audio in chat history
            st.session_state.chat_history[-1] = (prompt, answer, audio_b64)
